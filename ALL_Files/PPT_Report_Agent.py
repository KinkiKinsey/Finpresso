import os
import re
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_investment_presentation(json_file_path):
    """
    Creates a PowerPoint presentation based on investment data from a JSON file.
    
    Args:
        json_file_path (str): Path to the JSON file containing investment data
    
    Returns:
        str: Path to the created PowerPoint file
    """
    # Load the investment data from JSON
    with open(json_file_path, 'r') as file:
        data = json.load(file)
    
    # Extract ticker symbol
    ticker = data.get('Ticker', 'UNKNOWN')
    
    # Extract the Investment_Mindmap
    investment_mindmap = data.get('Investment_Mindmap', '')
    
    # Process the investment mindmap to extract sections
    three_key_takeaways = extract_key_takeaways(investment_mindmap)
    financial_situation_prospects = extract_financial_situation(investment_mindmap)
    market_catalysts = extract_catalysts(investment_mindmap)
    stock_price_volatility = extract_price_volatility(investment_mindmap)
    investment_recommendation = extract_investment_recommendation(investment_mindmap)
    
    # Create a new PowerPoint presentation
    ppt = Presentation()
    
    # Create Cover Slide
    create_cover_slide(ppt, ticker)
    
    # Create Key Takeaways Slide
    create_key_takeaways_slide(ppt, three_key_takeaways)
    
    # Create Financial Situation Slide
    create_financial_situation_slide(ppt, financial_situation_prospects)
    
    # Create Catalysts Slide
    create_catalysts_slide(ppt, market_catalysts)
    
    # Create Price & Volatility Slide
    create_price_volatility_slide(ppt, stock_price_volatility)
    
    # Create Investment Recommendation Slide
    create_investment_recommendation_slide(ppt, investment_recommendation)
    
    # Save the presentation
    ppt_file_path = os.path.join(os.getcwd(), f"{ticker}_financial_report.pptx")
    ppt.save(ppt_file_path)
    
    return ppt_file_path

def extract_key_takeaways(mindmap):
    """
    Extract key takeaways from the investment mindmap.
    
    Args:
        mindmap (str): The investment mindmap text
        
    Returns:
        dict: Dictionary of key takeaways
    """
    takeaways = {}
    
    # Try to find the company analysis section which typically contains key points
    company_section_match = re.search(r'COMPANY ANALYSIS.*?:(.*?)(?=\n\nPRICE ANALYSIS|\Z)', mindmap, re.DOTALL)
    if company_section_match:
        company_section = company_section_match.group(1).strip()
        
        # Extract 3 key points from the company section
        sentences = re.split(r'(?<=[.!?])\s+', company_section)
        sentences = [s.strip() for s in sentences if s.strip()]
        
        # Create up to 3 takeaways
        for i, sentence in enumerate(sentences[:3]):
            takeaways[f"Key Point {i+1}"] = sentence
            
    # If we couldn't find enough takeaways, add the recommendation
    if len(takeaways) < 3:
        thesis_match = re.search(r'INVESTMENT THESIS.*?:(.*?)(?=\n\n|\Z)', mindmap, re.DOTALL)
        if thesis_match:
            thesis = thesis_match.group(1).strip()
            takeaways["Recommendation"] = thesis
            
    # If still not enough, add from conclusion
    if len(takeaways) < 3:
        conclusion_match = re.search(r'CONCLUSION:(.*?)(?=\n\n|\Z)', mindmap, re.DOTALL)
        if conclusion_match:
            conclusion = conclusion_match.group(1).strip()
            takeaways["Conclusion"] = conclusion
    
    # If we still don't have enough, add placeholders
    for i in range(len(takeaways), 3):
        takeaways[f"Key Point {i+1}"] = "No additional key points available."
    
    return takeaways

def extract_financial_situation(mindmap):
    """
    Extract financial situation from the investment mindmap.
    
    Args:
        mindmap (str): The investment mindmap text
        
    Returns:
        dict: Dictionary with financial situation data
    """
    financial_data = {}
    
    # Try to find the macroeconomic environment section
    macro_section_match = re.search(r'MACROECONOMIC ENVIRONMENT:(.*?)(?=\n\nCOMPANY ANALYSIS|\Z)', mindmap, re.DOTALL)
    if macro_section_match:
        financial_data["Current Financial Situation"] = macro_section_match.group(1).strip()
    else:
        financial_data["Current Financial Situation"] = "No macroeconomic data available."
    
    # Try to find the company analysis section for future prospects
    company_section_match = re.search(r'COMPANY ANALYSIS.*?:(.*?)(?=\n\nPRICE ANALYSIS|\Z)', mindmap, re.DOTALL)
    if company_section_match:
        financial_data["Future Prospects"] = company_section_match.group(1).strip()
    else:
        financial_data["Future Prospects"] = "No company analysis data available."
    
    return financial_data

def create_cover_slide(ppt, ticker):
    """Creates the cover slide with the ticker symbol."""
    slide_layout = ppt.slide_layouts[5]  # Blank slide
    slide = ppt.slides.add_slide(slide_layout)
    
    # Set black background
    slide_background = slide.background
    fill = slide_background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)  # Black background
    
    # Add "AI Report of" text
    title_text = "AI Report of"
    title_box = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1))
    title_tf = title_box.text_frame
    title_tf.text = title_text
    
    title_p = title_tf.paragraphs[0]
    title_p.font.size = Pt(28)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(255, 255, 255)  # White text
    title_p.font.name = "Georgia"
    title_p.alignment = PP_ALIGN.CENTER
    
    # Add ticker symbol
    ticker_box = slide.shapes.add_textbox(Inches(2), Inches(4), Inches(6), Inches(2))
    ticker_tf = ticker_box.text_frame
    ticker_tf.text = ticker
    
    ticker_p = ticker_tf.paragraphs[0]
    ticker_p.font.size = Pt(72)
    ticker_p.font.bold = True
    ticker_p.font.color.rgb = RGBColor(255, 255, 255)  # White text
    ticker_p.font.name = "Arial"
    ticker_p.alignment = PP_ALIGN.CENTER

def create_key_takeaways_slide(ppt, takeaways):
    """Creates the key takeaways slide."""
    slide_layout = ppt.slide_layouts[5]  # Blank slide
    slide = ppt.slides.add_slide(slide_layout)
    
    # Set black background
    slide_background = slide.background
    fill = slide_background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)  # Black background
    
    # Add title
    title = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(8), Inches(1))
    title_tf = title.text_frame
    title_tf.text = "Three Key Takeaways⚠️"
    title_tf.paragraphs[0].font.size = Pt(30)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White text
    title_tf.paragraphs[0].font.name = "Georgia"
    title_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Define text layout
    text_left_margin = Inches(1)
    text_width = Inches(8)
    box_top_positions = [Inches(1.5), Inches(3.5), Inches(5.5)]
    
    # Font styles
    SUBTITLE_FONT = "Arial"
    CONTENT_FONT = "Consolas"
    TEXT_COLOR = RGBColor(255, 255, 255)
    
    # Add takeaways
    for i, (subtitle_text, description_text) in enumerate(takeaways.items()):
        if i >= 3:  # Limit to 3 takeaways
            break
            
        box_top = box_top_positions[i]
        
        # Add subtitle
        subtitle_box = slide.shapes.add_textbox(text_left_margin, box_top, text_width, Inches(0.5))
        subtitle_tf = subtitle_box.text_frame
        subtitle_tf.text = subtitle_text
        subtitle_tf.paragraphs[0].font.size = Pt(18)
        subtitle_tf.paragraphs[0].font.bold = True
        subtitle_tf.paragraphs[0].font.color.rgb = TEXT_COLOR
        subtitle_tf.paragraphs[0].font.name = SUBTITLE_FONT
        
        # Add description
        desc_box = slide.shapes.add_textbox(text_left_margin, box_top + Inches(0.5), text_width, Inches(1))
        desc_tf = desc_box.text_frame
        desc_tf.text = description_text
        desc_tf.word_wrap = True
        desc_tf.paragraphs[0].font.size = Pt(12)
        desc_tf.paragraphs[0].font.color.rgb = TEXT_COLOR
        desc_tf.paragraphs[0].font.name = CONTENT_FONT

def create_financial_situation_slide(ppt, financial_data):
    """Creates the financial situation slide."""
    slide_layout = ppt.slide_layouts[5]  # Blank slide
    slide = ppt.slides.add_slide(slide_layout)
    
    # Set black background
    slide_background = slide.background
    fill = slide_background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)  # Black background
    
    # Add title
    title = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(8), Inches(1))
    title_tf = title.text_frame
    title_tf.text = "Investment Environment and Future Prospects📊"
    title_tf.paragraphs[0].font.size = Pt(30)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White text
    title_tf.paragraphs[0].font.name = "Georgia"
    title_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Define layout
    left_x = Inches(1)
    right_x = Inches(5)
    text_width = Inches(4)
    box_top = Inches(1.5)
    
    # Font styles
    SUBTITLE_FONT = "Arial"
    CONTENT_FONT = "Consolas"
    TEXT_COLOR = RGBColor(255, 255, 255)
    
    # Get data for left and right sides
    keys = list(financial_data.keys())
    values = list(financial_data.values())
    
    # Add left side (Current Financial Situation)
    if len(keys) > 0:
        subtitle1 = slide.shapes.add_textbox(left_x, box_top, text_width, Inches(0.5))
        subtitle1_tf = subtitle1.text_frame
        subtitle1_tf.text = keys[0]
        subtitle1_tf.paragraphs[0].font.size = Pt(18)
        subtitle1_tf.paragraphs[0].font.bold = True
        subtitle1_tf.paragraphs[0].font.color.rgb = TEXT_COLOR
        subtitle1_tf.paragraphs[0].font.name = SUBTITLE_FONT
        
        desc1 = slide.shapes.add_textbox(left_x, box_top + Inches(0.5), text_width, Inches(2))
        desc1_tf = desc1.text_frame
        desc1_tf.text = values[0]
        desc1_tf.word_wrap = True
        desc1_tf.paragraphs[0].font.size = Pt(12)
        desc1_tf.paragraphs[0].font.color.rgb = TEXT_COLOR
        desc1_tf.paragraphs[0].font.name = CONTENT_FONT
    
    # Add right side (Future Prospects)
    if len(keys) > 1:
        subtitle2 = slide.shapes.add_textbox(right_x, box_top, text_width, Inches(0.5))
        subtitle2_tf = subtitle2.text_frame
        subtitle2_tf.text = keys[1]
        subtitle2_tf.paragraphs[0].font.size = Pt(18)
        subtitle2_tf.paragraphs[0].font.bold = True
        subtitle2_tf.paragraphs[0].font.color.rgb = TEXT_COLOR
        subtitle2_tf.paragraphs[0].font.name = SUBTITLE_FONT
        
        desc2 = slide.shapes.add_textbox(right_x, box_top + Inches(0.5), text_width, Inches(2))
        desc2_tf = desc2.text_frame
        desc2_tf.text = values[1]
        desc2_tf.word_wrap = True
        desc2_tf.paragraphs[0].font.size = Pt(12)
        desc2_tf.paragraphs[0].font.color.rgb = TEXT_COLOR
        desc2_tf.paragraphs[0].font.name = CONTENT_FONT

def create_catalysts_slide(ppt, catalysts):
    """Creates the market catalysts slide."""
    slide_layout = ppt.slide_layouts[5]  # Blank slide
    slide = ppt.slides.add_slide(slide_layout)
    
    # Get slide dimensions
    slide_width = ppt.slide_width
    slide_height = ppt.slide_height
    
    # Set black background
    slide_background = slide.background
    fill = slide_background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)  # Black background
    
    # Add title
    title = slide.shapes.add_textbox(Inches(0), Inches(0.3), slide_width, Inches(1))
    title_tf = title.text_frame
    title_tf.text = "Catalyst⏳"
    title_tf.paragraphs[0].font.size = Pt(40)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White text
    title_tf.paragraphs[0].font.name = "Georgia"
    title_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Define text layout
    text_left_margin = Inches(1)
    text_width = slide_width - Inches(2)
    box_top = Inches(1.5)
    
    # Define font scaling
    num_catalysts = len(catalysts)
    base_font_size = 16
    
    font_size = base_font_size
    if num_catalysts > 3:
        font_size = base_font_size - 4
    if num_catalysts > 6:
        font_size = base_font_size - 6.5
    if num_catalysts > 9:
        font_size = base_font_size - 9
    
    # Add catalyst items
    current_top = box_top
    
    for i, (key, value) in enumerate(catalysts.items()):
        # Create text box
        text_box = slide.shapes.add_textbox(text_left_margin, current_top, text_width, Inches(1))
        text_tf = text_box.text_frame
        text_tf.word_wrap = True
        text_tf.clear()
        
        # Create paragraph
        p = text_tf.add_paragraph()
        
        # Add bold key
        run_bold = p.add_run()
        run_bold.text = f"{key}: "
        run_bold.font.bold = True
        run_bold.font.size = Pt(font_size)
        run_bold.font.color.rgb = RGBColor(255, 255, 255)
        run_bold.font.name = "Arial"
        
        # Add regular value
        run_regular = p.add_run()
        run_regular.text = value
        run_regular.font.size = Pt(font_size * 0.8)
        run_regular.font.bold = False
        run_regular.font.color.rgb = RGBColor(255, 255, 255)
        run_regular.font.name = "Consolas"
        
        # Move to next position
        current_top += Inches(1.2)

def create_price_volatility_slide(ppt, price_volatility):
    """Creates the price and volatility analysis slide."""
    slide_layout = ppt.slide_layouts[5]  # Blank slide
    slide = ppt.slides.add_slide(slide_layout)
    
    # Set black background
    slide_background = slide.background
    fill = slide_background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)  # Black background
    
    # Add title
    title = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(8), Inches(1))
    title_tf = title.text_frame
    title_tf.text = "Stock Price & Volatility Analysis📈"
    title_tf.paragraphs[0].font.size = Pt(30)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White text
    title_tf.paragraphs[0].font.name = "Georgia"
    title_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Define layout
    left_x = Inches(0.5)
    right_x = Inches(5.5)
    bottom_x = Inches(0.5)
    
    upper_text_width = Inches(4)
    bottom_text_width = Inches(9)
    
    upper_box_top = Inches(1.5)
    bottom_box_top = Inches(4.5)
    
    # Font styles
    SUBTITLE_FONT = "Arial"
    CONTENT_FONT = "Consolas"
    TEXT_COLOR = RGBColor(255, 255, 255)
    
    # Get data
    keys = list(price_volatility.keys())
    values = list(price_volatility.values())
    
    # Add upper-left section (Stock Price Analysis)
    if len(keys) > 0:
        subtitle1 = slide.shapes.add_textbox(left_x, upper_box_top, upper_text_width, Inches(0.5))
        subtitle1_tf = subtitle1.text_frame
        subtitle1_tf.text = keys[0]
        subtitle1_tf.paragraphs[0].font.size = Pt(16)
        subtitle1_tf.paragraphs[0].font.bold = True
        subtitle1_tf.paragraphs[0].font.color.rgb = TEXT_COLOR
        subtitle1_tf.paragraphs[0].font.name = SUBTITLE_FONT
        
        desc1 = slide.shapes.add_textbox(left_x, upper_box_top + Inches(0.5), upper_text_width, Inches(2))
        desc1_tf = desc1.text_frame
        desc1_tf.text = values[0]
        desc1_tf.word_wrap = True
        desc1_tf.paragraphs[0].font.size = Pt(14)
        desc1_tf.paragraphs[0].font.color.rgb = TEXT_COLOR
        desc1_tf.paragraphs[0].font.name = CONTENT_FONT
    
    # Add upper-right section (Volatility Analysis)
    if len(keys) > 1:
        subtitle2 = slide.shapes.add_textbox(right_x, upper_box_top, upper_text_width, Inches(0.5))
        subtitle2_tf = subtitle2.text_frame
        subtitle2_tf.text = keys[1]
        subtitle2_tf.paragraphs[0].font.size = Pt(16)
        subtitle2_tf.paragraphs[0].font.bold = True
        subtitle2_tf.paragraphs[0].font.color.rgb = TEXT_COLOR
        subtitle2_tf.paragraphs[0].font.name = SUBTITLE_FONT
        
        desc2 = slide.shapes.add_textbox(right_x, upper_box_top + Inches(0.5), upper_text_width, Inches(2))
        desc2_tf = desc2.text_frame
        desc2_tf.text = values[1]
        desc2_tf.word_wrap = True
        desc2_tf.paragraphs[0].font.size = Pt(14)
        desc2_tf.paragraphs[0].font.color.rgb = TEXT_COLOR
        desc2_tf.paragraphs[0].font.name = CONTENT_FONT
    
    # Add bottom section (What They Reflect)
    if len(keys) > 2:
        subtitle3 = slide.shapes.add_textbox(bottom_x, bottom_box_top, bottom_text_width, Inches(0.5))
        subtitle3_tf = subtitle3.text_frame
        subtitle3_tf.text = keys[2]
        subtitle3_tf.paragraphs[0].font.size = Pt(16)
        subtitle3_tf.paragraphs[0].font.bold = True
        subtitle3_tf.paragraphs[0].font.color.rgb = TEXT_COLOR
        subtitle3_tf.paragraphs[0].font.name = SUBTITLE_FONT
        
        desc3 = slide.shapes.add_textbox(bottom_x, bottom_box_top + Inches(0.5), bottom_text_width, Inches(2))
        desc3_tf = desc3.text_frame
        desc3_tf.text = values[2]
        desc3_tf.word_wrap = True
        desc3_tf.paragraphs[0].font.size = Pt(14)
        desc3_tf.paragraphs[0].font.color.rgb = TEXT_COLOR
        desc3_tf.paragraphs[0].font.name = CONTENT_FONT

def create_investment_recommendation_slide(ppt, recommendations):
    """Creates the investment recommendation slide."""
    slide_layout = ppt.slide_layouts[5]  # Blank slide
    slide = ppt.slides.add_slide(slide_layout)
    
    # Set black background
    slide_background = slide.background
    fill = slide_background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)  # Black background
    
    # Add title
    title = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(8), Inches(1))
    title_tf = title.text_frame
    title_tf.text = "Investment Recommendation💰"
    title_tf.paragraphs[0].font.size = Pt(30)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White text
    title_tf.paragraphs[0].font.name = "Georgia"
    title_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Define layout
    text_left_margin = Inches(1)
    text_width = Inches(8)
    
    # Font styles
    SUBTITLE_FONT = "Arial"
    CONTENT_FONT = "Consolas"
    TEXT_COLOR = RGBColor(255, 255, 255)
    
    # Start position
    current_top = Inches(2)
    
    # Add recommendations
    for i, (subtitle_text, description_text) in enumerate(recommendations.items()):
        # Estimate line count
        approx_lines = max(3, len(description_text) // 80)
        additional_spacing = Inches(0.2) * (approx_lines - 3)
        
        # Add subtitle
        subtitle = slide.shapes.add_textbox(text_left_margin, current_top, text_width, Inches(0.5))
        subtitle_tf = subtitle.text_frame
        subtitle_tf.text = subtitle_text
        subtitle_tf.paragraphs[0].font.size = Pt(16)
        subtitle_tf.paragraphs[0].font.bold = True
        subtitle_tf.paragraphs[0].font.color.rgb = TEXT_COLOR
        subtitle_tf.paragraphs[0].font.name = SUBTITLE_FONT
        
        # Add description
        desc = slide.shapes.add_textbox(text_left_margin, current_top + Inches(0.4), text_width, Inches(1))
        desc_tf = desc.text_frame
        desc_tf.text = description_text
        desc_tf.word_wrap = True
        desc_tf.paragraphs[0].font.size = Pt(12)
        desc_tf.paragraphs[0].font.color.rgb = TEXT_COLOR
        desc_tf.paragraphs[0].font.name = CONTENT_FONT
        
        # Move to next section
        current_top += Inches(1.2) + additional_spacing

def extract_catalysts(mindmap):
    """
    Extract market catalysts from the investment mindmap.
    
    Args:
        mindmap (str): The investment mindmap text
        
    Returns:
        dict: Dictionary of catalysts
    """
    catalysts = {}
    
    # Look for specific catalysts in the integrated investment strategy section
    strategy_match = re.search(r'INTEGRATED INVESTMENT STRATEGY:(.*?)(?=\n\nCONCLUSION|\Z)', mindmap, re.DOTALL)
    if strategy_match:
        strategy_section = strategy_match.group(1).strip()
        
        # Try to find bullet points or similar indicators of catalysts
        bullet_points = re.findall(r'[-•]\s*(.*?)(?=\n[-•]|\n\n|\Z)', strategy_section, re.DOTALL)
        if bullet_points:
            for i, point in enumerate(bullet_points[:5]):  # Limit to 5 catalysts
                catalysts[f"Catalyst {i+1}"] = point.strip()
        
        # If no bullet points, try to break by sentences
        if not catalysts:
            sentences = re.split(r'(?<=[.!?])\s+', strategy_section)
            sentences = [s.strip() for s in sentences if s.strip()]
            for i, sentence in enumerate(sentences[:5]):  # Limit to 5 catalysts
                catalysts[f"Catalyst {i+1}"] = sentence
    
    # If still no catalysts, try to extract from the company analysis
    if not catalysts:
        company_match = re.search(r'COMPANY ANALYSIS.*?:(.*?)(?=\n\nPRICE ANALYSIS|\Z)', mindmap, re.DOTALL)
        if company_match:
            company_section = company_match.group(1).strip()
            sentences = re.split(r'(?<=[.!?])\s+', company_section)
            sentences = [s.strip() for s in sentences if s.strip()]
            for i, sentence in enumerate(sentences[:5]):  # Limit to 5 catalysts
                catalysts[f"Catalyst {i+1}"] = sentence
    
    # If we still don't have catalysts, add placeholders
    if not catalysts:
        catalysts["Catalyst 1"] = "No specific catalysts identified."
        catalysts["Catalyst 2"] = "Refer to company analysis for potential growth drivers."
    
    return catalysts

def extract_price_volatility(mindmap):
    """
    Extract price and volatility analysis from the investment mindmap.
    
    Args:
        mindmap (str): The investment mindmap text
        
    Returns:
        dict: Dictionary with price and volatility data
    """
    price_volatility = {}
    
    # Try to find the price analysis section
    price_match = re.search(r'PRICE ANALYSIS:(.*?)(?=\n\nINTEGRATED INVESTMENT STRATEGY|\Z)', mindmap, re.DOTALL)
    if price_match:
        price_section = price_match.group(1).strip()
        
        # Extract risk/reward analysis
        risk_reward_match = re.search(r'Risk/Reward Analysis.*?(?=\n\n|\Z)', price_section, re.DOTALL)
        if risk_reward_match:
            price_volatility["Stock Price Analysis"] = risk_reward_match.group(0).strip()
        else:
            price_volatility["Stock Price Analysis"] = "Risk/reward analysis not available."
        
        # Extract moving averages analysis
        ma_match = re.search(r'(?:SMA|EMA).*?Analysis.*?(?=\n\n|\Z)', price_section, re.DOTALL)
        if ma_match:
            price_volatility["Volatility Analysis"] = ma_match.group(0).strip()
        else:
            price_volatility["Volatility Analysis"] = "Moving average analysis not available."
        
        # Extract MACD analysis
        macd_match = re.search(r'MACD ANALYSIS.*?(?=\n\n|\Z)', price_section, re.DOTALL)
        if macd_match:
            price_volatility["What They Reflect"] = macd_match.group(0).strip()
        else:
            price_volatility["What They Reflect"] = "MACD analysis not available."
    else:
        # Default values if no price analysis found
        price_volatility["Stock Price Analysis"] = "No stock price analysis available."
        price_volatility["Volatility Analysis"] = "No volatility analysis available."
        price_volatility["What They Reflect"] = "No technical analysis reflection available."
    
    return price_volatility

def extract_investment_recommendation(mindmap):
    """
    Extract investment recommendation from the investment mindmap.
    
    Args:
        mindmap (str): The investment mindmap text
        
    Returns:
        dict: Dictionary with investment recommendation data
    """
    recommendation = {}
    
    # Try to find the integrated investment strategy section
    strategy_match = re.search(r'INTEGRATED INVESTMENT STRATEGY:(.*?)(?=\n\nCONCLUSION|\Z)', mindmap, re.DOTALL)
    if strategy_match:
        strategy_section = strategy_match.group(1).strip()
        
        # Extract position recommendation
        position_match = re.search(r'recommend(?:ation|ed)?\s+(?:a|an)?\s*(\w+)\s+strategy', strategy_section, re.IGNORECASE)
        if position_match:
            recommendation["What position we should take"] = f"We recommend a {position_match.group(1).upper()} position."
        else:
            # Try alternate patterns
            alt_position_match = re.search(r'(BUY|SELL|HOLD|WAIT|SHORT|LONG)\s+(?:is\s+)?recommend', strategy_section, re.IGNORECASE)
            if alt_position_match:
                recommendation["What position we should take"] = f"We recommend a {alt_position_match.group(1).upper()} position."
            else:
                recommendation["What position we should take"] = "Position recommendation not clearly specified."
        
        # Extract timeline
        timeline_match = re.search(r'(short-term|medium-term|long-term)', strategy_section, re.IGNORECASE)
        if timeline_match:
            recommendation["Investment timeline"] = f"This is a {timeline_match.group(1).upper()} investment opportunity."
        else:
            recommendation["Investment timeline"] = "Investment timeline not specified."
        
        # Extract risk level
        risk_match = re.search(r'(low|moderate|high)\s+risk', strategy_section, re.IGNORECASE)
        if risk_match:
            recommendation["Risk factors"] = f"This investment carries {risk_match.group(1).upper()} risk."
        else:
            recommendation["Risk factors"] = "Risk level not specified."
        
        # Extract rationale
        rationale_match = re.search(r'rationale is:(.*?)(?=\n|For optimal entry|\Z)', strategy_section, re.DOTALL)
        if rationale_match:
            recommendation["Rationale"] = rationale_match.group(1).strip()
        else:
            recommendation["Rationale"] = "Investment rationale not explicitly stated."
    else:
        # Default values if no strategy section found
        recommendation["What position we should take"] = "No position recommendation available."
        recommendation["Investment timeline"] = "No timeline information available."
        recommendation["Risk factors"] = "No risk assessment available."
        recommendation["Rationale"] = "No investment rationale available."
    
    return recommendation

def main():
    """Main function to run the agent."""
    import argparse
    
    parser = argparse.ArgumentParser(description='Generate PowerPoint from investment data JSON')
    parser.add_argument('json_file', help='Path to the JSON file containing investment data')
    args = parser.parse_args()
    
    ppt_path = create_investment_presentation(args.json_file)
    print(f"📊 PowerPoint presentation saved at: {ppt_path}")

if __name__ == "__main__":
    main()
